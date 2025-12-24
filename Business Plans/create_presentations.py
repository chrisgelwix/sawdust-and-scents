"""
Script to create PowerPoint presentations from business plan documents
Requires: pip install python-pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

def create_business_plan_presentation():
    """Create PowerPoint presentation from business plan"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Title Slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    add_title_slide(slide, "Sawdust and Scents", "Business Plan & Financial Projections", "2024")
    
    # Executive Summary
    slide = prs.slides.add_slide(prs.slide_layouts[1])  # Title and Content
    slide.shapes.title.text = "Executive Summary"
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "E-commerce venture specializing in handcrafted woodworking products and artisanal candles"
    p = tf.add_paragraph()
    p.text = "• Target Market: Home decor enthusiasts, gift buyers, DIY consumers"
    p = tf.add_paragraph()
    p.text = "• Initial Investment: $150,000 - $200,000"
    p = tf.add_paragraph()
    p.text = "• Break-Even: Month 12-15"
    p = tf.add_paragraph()
    p.text = "• 3-Year ROI: 118-120%"
    
    # Company Overview
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Company Overview"
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Mission: Create beautiful, handcrafted products that bring warmth to homes"
    p = tf.add_paragraph()
    p.text = "• Manufacturing: In-house production"
    p = tf.add_paragraph()
    p.text = "• Sales: E-commerce platform"
    p = tf.add_paragraph()
    p.text = "• Products: Woodworking items & candles"
    
    # Market Analysis
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Market Analysis"
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Market Size"
    p = tf.add_paragraph()
    p.text = "• US Home Decor Market: $120+ billion (4.5% CAGR)"
    p = tf.add_paragraph()
    p.text = "• US Candle Market: $3.2 billion (3.5% CAGR)"
    p = tf.add_paragraph()
    p.text = "• E-commerce Growth: 15-20% annually"
    
    # Financial Highlights
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Financial Highlights"
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Key Metrics"
    p = tf.add_paragraph()
    p.text = "• Initial Investment: $150,000"
    p = tf.add_paragraph()
    p.text = "• Break-Even: 280 units/month ($14,560)"
    p = tf.add_paragraph()
    p.text = "• Gross Margin: 54%"
    p = tf.add_paragraph()
    p.text = "• 3-Year ROI: 118-120%"
    
    # Revenue Projections
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "3-Year Revenue Projections"
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Year 1: $140,920 revenue, ($19,500) net"
    p = tf.add_paragraph()
    p.text = "Year 2: $343,200 revenue, $78,000 profit (22.7% margin)"
    p = tf.add_paragraph()
    p.text = "Year 3: $468,000 revenue, $120,000 profit (25.6% margin)"
    
    # Break-Even Analysis
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Break-Even Analysis"
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Break-Even Calculation"
    p = tf.add_paragraph()
    p.text = "• Fixed Costs: $8,400/month"
    p = tf.add_paragraph()
    p.text = "• Contribution Margin: $30/unit"
    p = tf.add_paragraph()
    p.text = "• Break-Even: 280 units/month"
    p = tf.add_paragraph()
    p.text = "• Timeline: Month 12-15"
    
    # ROI Analysis
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Return on Investment"
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "3-Year ROI: 118-120%"
    p = tf.add_paragraph()
    p.text = "• Year 1: ($21,000) - startup phase"
    p = tf.add_paragraph()
    p.text = "• Year 2: $78,000 profit (52% ROI)"
    p = tf.add_paragraph()
    p.text = "• Year 3: $120,000 profit (80% ROI)"
    p = tf.add_paragraph()
    p.text = "• Payback Period: 2.5 years"
    
    # Use of Funds
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Use of Funds"
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "$150,000 Investment Breakdown"
    p = tf.add_paragraph()
    p.text = "• Equipment & Setup: $48,000 (32%)"
    p = tf.add_paragraph()
    p.text = "• E-commerce Platform: $30,000 (20%)"
    p = tf.add_paragraph()
    p.text = "• Initial Inventory: $20,000 (13%)"
    p = tf.add_paragraph()
    p.text = "• Marketing & Launch: $15,000 (10%)"
    p = tf.add_paragraph()
    p.text = "• Working Capital: $37,000+ (25%+)"
    
    # Risk & Mitigation
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Risk Analysis & Mitigation"
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Key Risks & Mitigation Strategies"
    p = tf.add_paragraph()
    p.text = "• Market Risk: Strong branding, SEO, unique products"
    p = tf.add_paragraph()
    p.text = "• Production Risk: Quality control, multiple suppliers"
    p = tf.add_paragraph()
    p.text = "• Financial Risk: Conservative projections, working capital"
    
    # Growth Strategy
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Growth Strategy"
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "3-Year Growth Plan"
    p = tf.add_paragraph()
    p.text = "Year 1: Establish brand, achieve break-even"
    p = tf.add_paragraph()
    p.text = "Year 2: Expand products, increase marketing, hire staff"
    p = tf.add_paragraph()
    p.text = "Year 3: Scale production, explore wholesale partnerships"
    
    # Conclusion
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Investment Opportunity"
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Why Invest in Sawdust and Scents?"
    p = tf.add_paragraph()
    p.text = "✓ Break-even within 12-15 months"
    p = tf.add_paragraph()
    p.text = "✓ Strong gross margins (54%+)"
    p = tf.add_paragraph()
    p.text = "✓ Scalable business model"
    p = tf.add_paragraph()
    p.text = "✓ Growing market demand"
    p = tf.add_paragraph()
    p.text = "✓ 3-Year ROI: 118-120%"
    
    # Thank You
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_slide(slide, "Thank You", "Questions & Discussion", "")
    
    prs.save('Business_Plan_Presentation.pptx')
    print("Created Business_Plan_Presentation.pptx")

def create_financial_model_presentation():
    """Create PowerPoint presentation from financial model"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Title Slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_slide(slide, "Sawdust and Scents", "Financial Model & Analysis", "2024")
    
    # Key Metrics
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Key Financial Metrics"
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Quick Reference"
    p = tf.add_paragraph()
    p.text = "• Initial Investment: $150,000"
    p = tf.add_paragraph()
    p.text = "• Monthly Break-Even: 280 units / $14,560"
    p = tf.add_paragraph()
    p.text = "• Average Gross Margin: 54%"
    p = tf.add_paragraph()
    p.text = "• Break-Even Timeline: Month 12-15"
    p = tf.add_paragraph()
    p.text = "• 3-Year ROI: 118%"
    p = tf.add_paragraph()
    p.text = "• Payback Period: 2.5 years"
    
    # Break-Even Analysis
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Break-Even Analysis"
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Calculation"
    p = tf.add_paragraph()
    p.text = "Fixed Costs: $8,400/month"
    p = tf.add_paragraph()
    p.text = "Average Selling Price: $52"
    p = tf.add_paragraph()
    p.text = "Variable Cost: $22"
    p = tf.add_paragraph()
    p.text = "Contribution Margin: $30/unit"
    p = tf.add_paragraph()
    p.text = "Break-Even Units: 280/month"
    p = tf.add_paragraph()
    p.text = "Break-Even Revenue: $14,560/month"
    
    # Cost Structure
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Cost Structure"
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Fixed Costs (Monthly): $8,400"
    p = tf.add_paragraph()
    p.text = "• Rent: $1,500"
    p = tf.add_paragraph()
    p.text = "• Utilities: $400"
    p = tf.add_paragraph()
    p.text = "• Salaries: $4,000"
    p = tf.add_paragraph()
    p.text = "• Marketing: $2,000"
    p = tf.add_paragraph()
    p.text = "• Insurance & Software: $500"
    p = tf.add_paragraph()
    p.text = "Variable Cost: $22/unit (materials, labor, shipping)"
    
    # Revenue Projections Year 1
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Year 1 Revenue Projections"
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Monthly Growth Trajectory"
    p = tf.add_paragraph()
    p.text = "Month 1: 50 units, $2,600 revenue"
    p = tf.add_paragraph()
    p.text = "Month 6: 180 units, $9,360 revenue"
    p = tf.add_paragraph()
    p.text = "Month 9: 300 units, $15,600 revenue (First Profit)"
    p = tf.add_paragraph()
    p.text = "Month 12: 450 units, $23,400 revenue"
    p = tf.add_paragraph()
    p.text = "Year 1 Total: 2,710 units, $140,920 revenue"
    
    # 3-Year Financial Summary
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "3-Year Financial Summary"
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Year 1: $140,920 revenue, ($19,500) net"
    p = tf.add_paragraph()
    p.text = "Year 2: $343,200 revenue, $78,000 profit (22.7% margin)"
    p = tf.add_paragraph()
    p.text = "Year 3: $468,000 revenue, $120,000 profit (25.6% margin)"
    p = tf.add_paragraph()
    p.text = "3-Year Cumulative: $952,120 revenue, $178,500 profit"
    
    # Products Made vs Sold
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Production Efficiency"
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Products Made vs. Sold"
    p = tf.add_paragraph()
    p.text = "Year 1: 3,000 produced, 2,710 sold (99.2% efficiency)"
    p = tf.add_paragraph()
    p.text = "Year 2: 7,000 produced, 6,600 sold (94.3% efficiency)"
    p = tf.add_paragraph()
    p.text = "Year 3: 9,500 produced, 9,000 sold (94.7% efficiency)"
    p = tf.add_paragraph()
    p.text = "Inventory Turnover: 18-25x annually"
    
    # ROI Analysis
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "ROI Analysis"
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Investment: $150,000"
    p = tf.add_paragraph()
    p.text = "Year 1: ($21,000) loss"
    p = tf.add_paragraph()
    p.text = "Year 2: $78,000 profit (52% ROI)"
    p = tf.add_paragraph()
    p.text = "Year 3: $120,000 profit (80% ROI)"
    p = tf.add_paragraph()
    p.text = "3-Year Total Return: $177,000"
    p = tf.add_paragraph()
    p.text = "3-Year ROI: 118%"
    p = tf.add_paragraph()
    p.text = "Payback Period: 2.5 years"
    
    # Cash Flow
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Cash Flow Projections"
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Year 1 Cash Flow"
    p = tf.add_paragraph()
    p.text = "Peak Cash Requirement: ~$190,000 (Month 9)"
    p = tf.add_paragraph()
    p.text = "Month 9: First positive cash flow"
    p = tf.add_paragraph()
    p.text = "Month 12: $2,850 positive cash flow"
    p = tf.add_paragraph()
    p.text = "Year 2+: Consistently positive cash flow"
    
    # Sensitivity Analysis
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Sensitivity Analysis"
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Scenario Analysis"
    p = tf.add_paragraph()
    p.text = "Base Case: Break-even Month 9, 118% ROI"
    p = tf.add_paragraph()
    p.text = "Conservative (-20% sales): Break-even Month 14, 8% ROI"
    p = tf.add_paragraph()
    p.text = "Optimistic (+20% sales): Break-even Month 7, 32% ROI"
    
    # KPIs
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Key Performance Indicators"
    content = slide.placeholders[1]
    tf = content.text_frame
    tf.text = "Financial KPIs"
    p = tf.add_paragraph()
    p.text = "Gross Margin: 57.7% (Target: >55%)"
    p = tf.add_paragraph()
    p.text = "Net Margin Year 3: 25.6% (Target: >20%)"
    p = tf.add_paragraph()
    p.text = "Operational KPIs"
    p = tf.add_paragraph()
    p.text = "Inventory Turnover: 18-25x (Target: >12x)"
    p = tf.add_paragraph()
    p.text = "Production Efficiency: 94-99% (Target: >95%)"
    
    # Thank You
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title_slide(slide, "Thank You", "Financial Analysis Complete", "")
    
    prs.save('Financial_Model_Presentation.pptx')
    print("Created Financial_Model_Presentation.pptx")

def add_title_slide(slide, title, subtitle, date):
    """Helper to add a title slide"""
    # Add title
    left = Inches(1)
    top = Inches(2.5)
    width = Inches(8)
    height = Inches(1)
    title_box = slide.shapes.add_textbox(left, top, width, height)
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.paragraphs[0].font.size = Pt(44)
    title_frame.paragraphs[0].font.bold = True
    title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Add subtitle
    top = Inches(4)
    subtitle_box = slide.shapes.add_textbox(left, top, width, height)
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = subtitle
    subtitle_frame.paragraphs[0].font.size = Pt(24)
    subtitle_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    
    # Add date if provided
    if date:
        top = Inches(5.5)
        date_box = slide.shapes.add_textbox(left, top, width, Inches(0.5))
        date_frame = date_box.text_frame
        date_frame.text = date
        date_frame.paragraphs[0].font.size = Pt(18)
        date_frame.paragraphs[0].alignment = PP_ALIGN.CENTER

if __name__ == "__main__":
    try:
        create_business_plan_presentation()
        create_financial_model_presentation()
        print("\nSuccessfully created both PowerPoint presentations!")
        print("   - Business_Plan_Presentation.pptx")
        print("   - Financial_Model_Presentation.pptx")
    except ImportError:
        print("Error: python-pptx not installed")
        print("   Please run: pip install python-pptx")
    except Exception as e:
        print(f"Error creating presentations: {e}")

